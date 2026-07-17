import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_ultra_premium_presentation():
    prs = Presentation()
    # 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # McKinsey/Vercel Inspired Ultra-Premium Color Palette
    COLOR_BG_DARK = RGBColor(10, 10, 15)       # Deep Obsidian Black #0a0a0f
    COLOR_BG_LIGHT = RGBColor(248, 250, 252)   # Slate-50 Canvas #f8fafc
    COLOR_CARD_BG = RGBColor(255, 255, 255)    # Pure White card #ffffff
    COLOR_INDIGO = RGBColor(79, 70, 229)       # Indigo Brand Accent #4f46e5
    COLOR_TEAL = RGBColor(14, 116, 144)        # Teal Brand Accent #0e7490
    COLOR_TEXT_PRIMARY = RGBColor(15, 23, 42)  # Slate-900 Primary #0f172a
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate-500 Muted #64748b
    COLOR_BORDER = RGBColor(226, 232, 240)     # Slate-200 Border #e2e8f0

    # Helper: Slide Background
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Premium Slide Header
    def add_premium_header(slide, title_text, category_text="PROJECT SPECIFICATION"):
        # Category label (uppercase small tracker)
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_right = cat_tf.margin_bottom = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category_text.upper()
        cat_p.font.name = 'Arial'
        cat_p.font.size = Pt(10)
        cat_p.font.bold = True
        cat_p.font.color.rgb = COLOR_INDIGO
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_right = title_tf.margin_bottom = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = 'Trebuchet MS'
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_TEXT_PRIMARY

        # Thin elegant separator line
        sep_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.02))
        sep_line.fill.solid()
        sep_line.fill.fore_color.rgb = COLOR_BORDER
        sep_line.line.fill.background()

    # --- SLIDE 1: Title Slide (Minimalist Obsidian theme) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG_DARK)

    # Elegant structural block graphic
    graphic = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.08), Inches(3.2))
    graphic.fill.solid()
    graphic.fill.fore_color.rgb = COLOR_INDIGO
    graphic.line.fill.background()

    # Title box
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_tag = tf.paragraphs[0]
    p_tag.text = "P R O M P T   E N G I N E E R I N G   C A S E   S T U D Y"
    p_tag.font.name = 'Arial'
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_TEXT_MUTED
    p_tag.space_after = Pt(14)

    p_title = tf.add_paragraph()
    p_title.text = "Syntactic Intelligence"
    p_title.font.name = 'Trebuchet MS'
    p_title.font.size = Pt(56)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    p_title.space_after = Pt(8)

    p_subtitle = tf.add_paragraph()
    p_subtitle.text = "Engineering Prompt Variations to Drive Multi-Functional AI Behaviors"
    p_subtitle.font.name = 'Arial'
    p_subtitle.font.size = Pt(20)
    p_subtitle.font.color.rgb = COLOR_BORDER
    p_subtitle.space_after = Pt(36)

    p_meta = tf.add_paragraph()
    p_meta.text = "System: Flask Framework | Core LLM: Llama-3.1-8B NIM API | Design: Flat Corporate Minimalist"
    p_meta.font.name = 'Arial'
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = COLOR_TEXT_MUTED

    # --- SLIDE 2: Project Overview & Strategy (Asymmetric layout) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG_LIGHT)
    add_premium_header(slide, "Overview & Prompt Strategy", "01. Introduction")

    # Column 1: Core Objectives (Left White Card)
    c1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.6), Inches(4.7))
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLOR_CARD_BG
    c1.line.color.rgb = COLOR_BORDER
    
    # Small top colored band for card premium aesthetic
    c1_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.6), Inches(0.08))
    c1_band.fill.solid()
    c1_band.fill.fore_color.rgb = COLOR_INDIGO
    c1_band.line.fill.background()

    tx1 = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.1))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Educational Focus"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    p.space_after = Pt(12)
    
    p_desc = tf1.add_paragraph()
    p_desc.text = "This project aims to enhance students' capabilities to formulate systemic prompts that inspire structured thinking, logical evaluation, and entrepreneurial creativity."
    p_desc.font.name = 'Calibri'
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = COLOR_TEXT_PRIMARY
    p_desc.line_spacing = 1.3
    p_desc.space_after = Pt(14)
    
    p_desc2 = tf1.add_paragraph()
    p_desc2.text = "By comparing response outputs side-by-side, developers learn to design constraints (length, specificity, output format) that yield optimal business and creative utility."
    p_desc2.font.name = 'Calibri'
    p_desc2.font.size = Pt(15)
    p_desc2.font.color.rgb = COLOR_TEXT_PRIMARY
    p_desc2.line_spacing = 1.3

    # Column 2: Key Outcomes (Right White Card)
    c2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.9), Inches(5.6), Inches(4.7))
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_CARD_BG
    c2.line.color.rgb = COLOR_BORDER
    
    c2_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.9), Inches(5.6), Inches(0.08))
    c2_band.fill.solid()
    c2_band.fill.fore_color.rgb = COLOR_TEAL
    c2_band.line.fill.background()

    tx2 = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.0), Inches(4.1))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "Expected Outcomes"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    p.space_after = Pt(14)

    bullets = [
        "A collection of highly structured prompt variations that yield unique, actionable AI outputs.",
        "Demonstrated understanding of tone, complexity, and format control through system instructions.",
        "An active feedback loop that stores and logs rating statistics to drive systematic prompt optimization.",
        "A highly intuitive, developer-centric single page web workspace."
    ]
    for bullet in bullets:
        bp = tf2.add_paragraph()
        bp.text = "• " + bullet
        bp.font.name = 'Calibri'
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT_PRIMARY
        bp.space_after = Pt(10)

    # --- SLIDE 3: System Architecture (Premium Modern Layout) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG_LIGHT)
    add_premium_header(slide, "System Architecture & Execution Flow", "02. Technology Stack")

    architectures = [
        ("01", "Backend Flow", "Flask Core", [
            "Handles endpoint routing securely.",
            "Assembles dynamic prompt messages.",
            "Pipeline endpoints: /, /generate, /feedback."
        ], Inches(0.8), COLOR_INDIGO),
        ("02", "AI Model Pipeline", "NVIDIA NIM", [
            "Model: Llama-3.1-8B-Instruct.",
            "Maintains 9 separate configurations.",
            "Requests-based chat completion endpoint."
        ], Inches(4.8), COLOR_TEAL),
        ("03", "Data Metrics", "CSV Database", [
            "Records logs in feedback_log.csv.",
            "Thread-locked structure to ensure safety.",
            "Dashboard: aggregates analytics at runtime."
        ], Inches(8.8), COLOR_TEXT_PRIMARY)
    ]

    for number, title, subtitle, bullets, left_pos, color in architectures:
        # Card Surface
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.9), Inches(3.7), Inches(4.7))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER
        
        # Indicator Number (Large, high-contrast)
        num_box = slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.0), Inches(1.5), Inches(0.8))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = number
        num_p.font.name = 'Trebuchet MS'
        num_p.font.size = Pt(40)
        num_p.font.bold = True
        num_p.font.color.rgb = color
        
        tx = slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.9), Inches(3.3), Inches(3.5))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Trebuchet MS'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_PRIMARY
        p_title.space_after = Pt(2)
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle.upper()
        p_sub.font.name = 'Arial'
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_TEXT_MUTED
        p_sub.space_after = Pt(12)
        
        for bullet in bullets:
            bp = tf.add_paragraph()
            bp.text = "• " + bullet
            bp.font.name = 'Calibri'
            bp.font.size = Pt(13)
            bp.font.color.rgb = COLOR_TEXT_PRIMARY
            bp.space_after = Pt(8)

    # --- SLIDES 4, 5, 6: Functions & Prompt Designs (Ultra-Premium Layout) ---
    def create_premium_function_slide(number_str, title_text, system_desc, prompts_list):
        f_slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(f_slide, COLOR_BG_LIGHT)
        add_premium_header(f_slide, title_text, f"03.{number_str} Function Details")
        
        # System instructions panel (Full-width light card)
        sys_bg = f_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.9))
        sys_bg.fill.solid()
        sys_bg.fill.fore_color.rgb = COLOR_CARD_BG
        sys_bg.line.color.rgb = COLOR_BORDER
        
        sys_tx = f_slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.3), Inches(0.7))
        sys_tf = sys_tx.text_frame
        sys_tf.word_wrap = True
        sys_p = sys_tf.paragraphs[0]
        sys_p.text = "SYSTEM ROLE & CONTEXT: " + system_desc
        sys_p.font.name = 'Arial'
        sys_p.font.size = Pt(11)
        sys_p.font.bold = True
        sys_p.font.color.rgb = COLOR_TEXT_MUTED
        sys_p.space_after = Pt(4)
        
        sys_sub = sys_tf.add_paragraph()
        sys_sub.text = "Provides specialized guidance to Llama-3.1 before formatting the user queries."
        sys_sub.font.name = 'Calibri'
        sys_sub.font.size = Pt(13)
        sys_sub.font.color.rgb = COLOR_TEXT_PRIMARY

        # 3 Prompt variation columns
        for i, (name, prompt_tmpl, design_intent) in enumerate(prompts_list):
            left_pos = Inches(0.8 + i * 4.0)
            card = f_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(2.9), Inches(3.7), Inches(3.7))
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_CARD_BG
            card.line.color.rgb = COLOR_BORDER
            
            # Colored top indicator
            card_accent = f_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(2.9), Inches(3.7), Inches(0.06))
            card_accent.fill.solid()
            card_accent.fill.fore_color.rgb = COLOR_INDIGO if i != 1 else COLOR_TEAL
            card_accent.line.fill.background()

            tx = f_slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(3.1), Inches(3.3), Inches(3.3))
            tf = tx.text_frame
            tf.word_wrap = True
            
            p_head = tf.paragraphs[0]
            p_head.text = name
            p_head.font.name = 'Trebuchet MS'
            p_head.font.size = Pt(18)
            p_head.font.bold = True
            p_head.font.color.rgb = COLOR_TEXT_PRIMARY
            p_head.space_after = Pt(8)
            
            p_lbl1 = tf.add_paragraph()
            p_lbl1.text = "PROMPT TEMPLATE"
            p_lbl1.font.name = 'Arial'
            p_lbl1.font.size = Pt(9)
            p_lbl1.font.bold = True
            p_lbl1.font.color.rgb = COLOR_TEXT_MUTED
            p_lbl1.space_after = Pt(2)
            
            p_tmpl = tf.add_paragraph()
            p_tmpl.text = f'"{prompt_tmpl}"'
            p_tmpl.font.name = 'Calibri'
            p_tmpl.font.size = Pt(13)
            p_tmpl.font.bold = True
            p_tmpl.font.color.rgb = COLOR_TEXT_PRIMARY
            p_tmpl.space_after = Pt(12)
            
            p_lbl2 = tf.add_paragraph()
            p_lbl2.text = "PROMPT STRATEGY"
            p_lbl2.font.name = 'Arial'
            p_lbl2.font.size = Pt(9)
            p_lbl2.font.bold = True
            p_lbl2.font.color.rgb = COLOR_TEXT_MUTED
            p_lbl2.space_after = Pt(2)
            
            p_strat = tf.add_paragraph()
            p_strat.text = design_intent
            p_strat.font.name = 'Calibri'
            p_strat.font.size = Pt(13)
            p_strat.font.color.rgb = COLOR_TEXT_PRIMARY

    # --- SLIDE 4: Function 1: Answer Questions ---
    create_premium_function_slide(
        "A",
        "1. Function: Answer Questions",
        "You are a knowledgeable AI assistant specialized in answering questions accurately. Provide clear, well-structured answers.",
        [
            ("Short Factual", "Answer briefly and factually: {input}", "Restricts response length. Eliminates preambles, focusing solely on core, direct facts."),
            ("Detailed Explainer", "Explain in detail with context and examples: {input}", "Instructs Llama-3.1 to generate fully fleshed-out contextual write-ups with use cases."),
            ("Bullet-Point Facts", "Give me 3 concise facts about: {input}", "Strict structural layout constraint. Limits the output to exactly 3 high-contrast facts.")
        ]
    )

    # --- SLIDE 5: Function 1 Screenshot (App Response Mockup) ---
    slide_qa_ss = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide_qa_ss, COLOR_BG_LIGHT)
    add_premium_header(slide_qa_ss, "Q&A System Response Interface Mockup", "03.A Response Interface")
    
    # Check if screenshot file exists
    if os.path.exists("qa_screen.png"):
        # Add screenshot centered on the right
        slide_qa_ss.shapes.add_picture("qa_screen.png", Inches(5.8), Inches(1.8), width=Inches(6.7))
        
        # Details on the left (Description Card)
        c_desc = slide_qa_ss.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(4.6), Inches(4.8))
        c_desc.fill.solid()
        c_desc.fill.fore_color.rgb = COLOR_CARD_BG
        c_desc.line.color.rgb = COLOR_BORDER
        
        tx = slide_qa_ss.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(4.2), Inches(4.2))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Interface Details"
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_INDIGO
        p.space_after = Pt(14)
        
        bullets = [
            "Query Input: User asks 'what is ai' under the Answer Questions workspace.",
            "Prompt Style: 'Bullet-Point Facts' dropdown is selected.",
            "AI Output: Llama-3.1 renders 3 concise facts structured under sub-headers.",
            "Interactive Action: Copy clipboard and helpful feedback controls log seamlessly."
        ]
        for bullet in bullets:
            bp = tf.add_paragraph()
            bp.text = "• " + bullet
            bp.font.name = 'Calibri'
            bp.font.size = Pt(13)
            bp.font.color.rgb = COLOR_TEXT_PRIMARY
            bp.space_after = Pt(10)
    else:
        # Fallback text
        tx = slide_qa_ss.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
        tx.text_frame.text = "[Screenshot 'qa_screen.png' not found in workspace]"

    # --- SLIDE 6: Function 2: Summarize Text ---
    create_premium_function_slide(
        "B",
        "2. Function: Summarize Text",
        "You are an expert text summarizer. Distill the provided content into its most essential points while preserving key meaning.",
        [
            ("One-Paragraph", "Summarize the following text in one paragraph: {input}", "Aggregates lengthy passages into a single cohesive synthesis block, optimized for quick scanning."),
            ("Key Points List", "Extract the main points of this text as a bulleted list: {input}", "Structures key topics as clear bullet lists, improving legibility of dense reports."),
            ("Executive Brief", "Provide a brief executive overview (2-3 sentences) of this document: {input}", "Generates hyper-condensed high-level summaries for fast management assessments.")
        ]
    )

    # --- SLIDE 7: Function 3: Generate Creative Content ---
    create_premium_function_slide(
        "C",
        "3. Function: Generate Creative Content",
        "You are a creative writing assistant with a flair for vivid imagery and engaging narratives. Produce original, compelling content.",
        [
            ("Short Story", "Write a short story (under 300 words) about: {input}", "Tests narrative development capabilities while enforcing a strict parameter length constraint."),
            ("Poem", "Write a poem in a creative tone about: {input}", "Directs the AI model to express metaphors, rhythm, and artistic structures for creative tasks."),
            ("Idea Generator", "Generate 3 creative concept ideas for a story/essay about: {input}", "Enables ideation brainstorming, producing three unique business/creative outlines.")
        ]
    )

    # --- SLIDE 8: Feedback Loop & Refinement (Consulting Style) ---
    slide_fb = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide_fb, COLOR_BG_LIGHT)
    add_premium_header(slide_fb, "Feedback Loop & Analytics Dashboard", "04. Optimization")

    # Left Card: Rating Logging
    l_card = slide_fb.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.6), Inches(4.7))
    l_card.fill.solid()
    l_card.fill.fore_color.rgb = COLOR_CARD_BG
    l_card.line.color.rgb = COLOR_BORDER
    
    l_band = slide_fb.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.6), Inches(0.08))
    l_band.fill.solid()
    l_band.fill.fore_color.rgb = COLOR_INDIGO
    l_band.line.fill.background()

    tx1 = slide_fb.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(5.0), Inches(4.1))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Interactive Feedback Mechanism"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    p.space_after = Pt(14)
    
    bullets = [
        "Inline Rating: 👍 Helpful / 👎 Not Helpful buttons appear dynamically immediately beneath each generated response block.",
        "Asynchronous AJAX: Submit ratings immediately without page reloading to maximize user convenience.",
        "Logged Columns: Records timestamp, active function, selected style variation, truncated input prompt, LLM response, and feedback rating."
    ]
    for bullet in bullets:
        bp = tf1.add_paragraph()
        bp.text = "• " + bullet
        bp.font.name = 'Calibri'
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT_PRIMARY
        bp.space_after = Pt(12)

    # Right Card: Prompt Evaluation
    r_card = slide_fb.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.9), Inches(5.6), Inches(4.7))
    r_card.fill.solid()
    r_card.fill.fore_color.rgb = COLOR_CARD_BG
    r_card.line.color.rgb = COLOR_BORDER
    
    r_band = slide_fb.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.9), Inches(5.6), Inches(0.08))
    r_band.fill.solid()
    r_band.fill.fore_color.rgb = COLOR_TEAL
    r_band.line.fill.background()

    tx2 = slide_fb.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.0), Inches(4.1))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "Evaluation & Prompt Iteration"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    p.space_after = Pt(14)
    
    bullets = [
        "Aggregate Dashboard: The /feedback-summary route calculates approval ratings across the 9 prompts dynamically.",
        "Empirical Alignment: Visual statistics highlight which prompt styles meet target audience needs.",
        "Prompt Tuning cycle: Provides the datasets required for students to systematically refine system templates."
    ]
    for bullet in bullets:
        bp = tf2.add_paragraph()
        bp.text = "• " + bullet
        bp.font.name = 'Calibri'
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT_PRIMARY
        bp.space_after = Pt(12)

    # --- SLIDE 9: Analytics Screen Screenshot ---
    slide_an_ss = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide_an_ss, COLOR_BG_LIGHT)
    add_premium_header(slide_an_ss, "Feedback Summary & Analytics Dashboard Mockup", "04.B Summary Interface")
    
    # Check if analytics screenshot exists
    if os.path.exists("analytics_screen.png"):
        # Add screenshot centered on the right
        slide_an_ss.shapes.add_picture("analytics_screen.png", Inches(5.8), Inches(1.8), width=Inches(6.7))
        
        # Details on the left (Description Card)
        c_desc = slide_an_ss.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(4.6), Inches(4.8))
        c_desc.fill.solid()
        c_desc.fill.fore_color.rgb = COLOR_CARD_BG
        c_desc.line.color.rgb = COLOR_BORDER
        
        tx = slide_an_ss.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(4.2), Inches(4.2))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Analytics Panel"
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_INDIGO
        p.space_after = Pt(14)
        
        bullets = [
            "Approval Aggregates: Displays total user ratings, helpful count, and satisfaction rates.",
            "Visual Metrics: A dynamic progress bar indicating satisfied vs dissatisfied ratios.",
            "Sub-Breakdowns: Tabulates thumbs up/down results for all 9 prompt styles individually.",
            "Real-time Updates: Refreshes dynamically as user feedback is logged to CSV."
        ]
        for bullet in bullets:
            bp = tf.add_paragraph()
            bp.text = "• " + bullet
            bp.font.name = 'Calibri'
            bp.font.size = Pt(13)
            bp.font.color.rgb = COLOR_TEXT_PRIMARY
            bp.space_after = Pt(10)
    else:
        # Fallback text
        tx = slide_an_ss.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
        tx.text_frame.text = "[Screenshot 'analytics_screen.png' not found in workspace]"

    output_path = "AI_Assistant_User_Guide_Final.pptx"
    prs.save(output_path)
    print(f"Ultra-Premium Presentation with screenshots saved to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_ultra_premium_presentation()
